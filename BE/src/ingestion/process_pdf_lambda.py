import io
import json
import os
import re
import time
import uuid
import zipfile
from datetime import datetime, timezone
from pathlib import PurePosixPath
from urllib.parse import unquote_plus
from xml.etree import ElementTree as ET

import boto3
from botocore.exceptions import ClientError

try:
    import pdfplumber
except Exception:  # pragma: no cover
    pdfplumber = None

try:
    from pypdf import PdfReader
except Exception:  # pragma: no cover
    PdfReader = None

try:
    from docx import Document as DocxDocument
except Exception:  # pragma: no cover
    DocxDocument = None

try:
    from pptx import Presentation
except Exception:  # pragma: no cover
    Presentation = None


AWS_REGION = os.environ.get("AWS_REGION", "ap-southeast-1")
DOCUMENTS_TABLE = os.environ.get("DOCUMENTS_TABLE", "StudyBotDocuments")
UPLOADS_BUCKET_NAME = os.environ.get("UPLOADS_BUCKET_NAME", "")
BEDROCK_KNOWLEDGE_BASE_ID = os.environ.get("BEDROCK_KNOWLEDGE_BASE_ID", "")
BEDROCK_DATA_SOURCE_ID = os.environ.get("BEDROCK_DATA_SOURCE_ID", "")
KB_PROCESSED_PREFIX = os.environ.get("KB_PROCESSED_PREFIX", "processed")
USE_TEXTRACT_FALLBACK = os.environ.get("USE_TEXTRACT_FALLBACK", "true").lower() == "true"
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".md", ".markdown", ".txt", ".pptx", ".vtt"}


s3 = boto3.client("s3", region_name=AWS_REGION)
ddb = boto3.resource("dynamodb", region_name=AWS_REGION).Table(DOCUMENTS_TABLE)
textract = boto3.client("textract", region_name=AWS_REGION)
bedrock_agent = boto3.client("bedrock-agent", region_name=AWS_REGION)


def now_iso():
    return datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")


def parse_object_key(key):
    # New upload key: raw/{user_id}/{session_id}/{doc_id}/{filename}
    match = re.match(r"^raw/(?P<user_id>[^/]+)/(?P<session_id>[^/]+)/(?P<doc_id>[^/]+)/(?P<filename>.+)$", key)
    if match:
        return match.groupdict()
    # Legacy upload key: documents/raw/{user_id}/{doc_id}/{filename}
    match = re.match(r"^documents/raw/(?P<user_id>[^/]+)/(?P<doc_id>[^/]+)/(?P<filename>.+)$", key)
    if match:
        parsed = match.groupdict()
        parsed["session_id"] = ""
        return parsed
    return None


def update_doc(user_id, doc_id, **attrs):
    item = ddb.get_item(Key={"PK": f"USER#{user_id}", "SK": f"DOC#{doc_id}"}).get("Item")
    if not item:
        item = {"PK": f"USER#{user_id}", "SK": f"DOC#{doc_id}", "doc_id": doc_id}
    item.update(attrs)
    ddb.put_item(Item=item)


def extract_with_pdfplumber(pdf_bytes):
    if not pdfplumber:
        return {"ok": False, "reason": "pdfplumber_not_installed"}

    text_parts = []
    page_count = 0
    non_empty_pages = 0
    table_rows = 0
    table_cells = 0

    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        page_count = len(pdf.pages)
        for page in pdf.pages:
            text = (page.extract_text() or "").strip()
            if text:
                non_empty_pages += 1
                text_parts.append(text)
            tables = page.extract_tables() or []
            for table in tables:
                table_rows += len(table or [])
                for row in table or []:
                    for cell in row or []:
                        if str(cell or "").strip():
                            table_cells += 1

    full_text = "\n\n".join(text_parts).strip()
    density = len(full_text) / max(page_count, 1)
    non_empty_ratio = non_empty_pages / max(page_count, 1)
    table_quality = table_cells / max(table_rows, 1) if table_rows else 0

    scanned_like = non_empty_ratio < 0.35 or density < 90
    table_poor = table_rows > 0 and table_quality < 1.5
    ok = bool(full_text) and not scanned_like and not table_poor

    return {
        "ok": ok,
        "text": full_text,
        "page_count": page_count,
        "density": density,
        "table_quality": table_quality,
        "scanned_like": scanned_like,
        "table_poor": table_poor,
    }


def extract_with_pypdf(pdf_bytes):
    if not PdfReader:
        return {"ok": False, "reason": "pypdf_not_installed"}

    reader = PdfReader(io.BytesIO(pdf_bytes))
    pages = []
    for page in reader.pages:
        pages.append((page.extract_text() or "").strip())
    text = "\n\n".join([p for p in pages if p]).strip()
    density = len(text) / max(len(reader.pages), 1)
    return {"ok": bool(text and density >= 90), "text": text, "page_count": len(reader.pages), "density": density}


def decode_text_bytes(file_bytes):
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="replace")


def extract_plain_text(file_bytes):
    text = decode_text_bytes(file_bytes).strip()
    return {"ok": bool(text), "text": text}


def is_vtt_timestamp(value):
    return bool(
        re.match(
            r"^\d{2}:\d{2}:\d{2}\.\d{3}\s+-->\s+\d{2}:\d{2}:\d{2}\.\d{3}",
            value or "",
        )
    )


def strip_vtt_markup(value):
    text = re.sub(r"<[^>]+>", "", value or "")
    text = re.sub(r"&nbsp;", " ", text)
    text = re.sub(r"&amp;", "&", text)
    text = re.sub(r"&lt;", "<", text)
    text = re.sub(r"&gt;", ">", text)
    return " ".join(text.split()).strip()


def extract_vtt_text(file_bytes):
    raw_text = decode_text_bytes(file_bytes)
    lines = raw_text.replace("\ufeff", "").splitlines()

    cues = []
    pending_time = ""
    pending_text = []
    seen_consecutive = set()

    def flush():
        nonlocal pending_time, pending_text, seen_consecutive
        if not pending_time:
            pending_text = []
            return

        cue_text = " ".join(item for item in pending_text if item).strip()
        if cue_text:
            start_time = pending_time.split("-->", 1)[0].strip()
            key = cue_text.lower()
            if key not in seen_consecutive:
                cues.append(f"[{start_time}] {cue_text}")
            seen_consecutive = {key}

        pending_time = ""
        pending_text = []

    skip_block = False
    for raw_line in lines:
        line = raw_line.strip()
        if not line:
            flush()
            skip_block = False
            continue

        if line.upper() == "WEBVTT":
            continue
        if line.startswith(("NOTE", "STYLE", "REGION")):
            skip_block = True
            continue
        if skip_block:
            continue

        if is_vtt_timestamp(line):
            flush()
            pending_time = line
            continue

        if not pending_time:
            # Cue identifier.
            continue

        cleaned = strip_vtt_markup(line)
        if cleaned:
            pending_text.append(cleaned)

    flush()

    text = "\n".join(cues).strip()
    return {"ok": bool(text), "text": text, "cue_count": len(cues), "mode": "webvtt"}


def iter_docx_table_text(table):
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
        if cells:
            yield " | ".join(cells)


def extract_with_python_docx(file_bytes):
    if not DocxDocument:
        return {"ok": False, "reason": "python_docx_not_installed"}

    try:
        document = DocxDocument(io.BytesIO(file_bytes))
    except Exception:
        return {"ok": False, "reason": "python_docx_parse_failed"}

    parts = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table in document.tables:
        parts.extend(iter_docx_table_text(table))

    text = "\n\n".join(parts).strip()
    return {"ok": bool(text), "text": text}


def local_name(tag):
    return tag.rsplit("}", 1)[-1]


def xml_text_blocks(xml_bytes, paragraph_tag="p"):
    root = ET.fromstring(xml_bytes)
    blocks = []
    for node in root.iter():
        if local_name(node.tag) != paragraph_tag:
            continue
        texts = []
        for child in node.iter():
            if local_name(child.tag) == "t" and child.text:
                texts.append(child.text)
        block = "".join(texts).strip()
        if block:
            blocks.append(block)
    return blocks


def extract_docx_from_zip(file_bytes):
    parts = []
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as archive:
            for name in (
                "word/document.xml",
                "word/footnotes.xml",
                "word/endnotes.xml",
            ):
                if name in archive.namelist():
                    parts.extend(xml_text_blocks(archive.read(name)))
    except (zipfile.BadZipFile, ET.ParseError):
        return {"ok": False, "reason": "docx_zip_parse_failed"}

    text = "\n\n".join(parts).strip()
    return {"ok": bool(text), "text": text}


def extract_docx_text(file_bytes):
    extraction = extract_with_python_docx(file_bytes)
    if extraction.get("ok"):
        extraction["mode"] = "python-docx"
        return extraction

    extraction = extract_docx_from_zip(file_bytes)
    if extraction.get("ok"):
        extraction["mode"] = "docx-xml"
    return extraction


def iter_pptx_shape_text(shape):
    if getattr(shape, "has_text_frame", False) and shape.text_frame:
        text = shape.text_frame.text.strip()
        if text:
            yield text

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                yield " | ".join(cells)

    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            yield from iter_pptx_shape_text(child)


def extract_with_python_pptx(file_bytes):
    if not Presentation:
        return {"ok": False, "reason": "python_pptx_not_installed"}

    try:
        presentation = Presentation(io.BytesIO(file_bytes))
    except Exception:
        return {"ok": False, "reason": "python_pptx_parse_failed"}

    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            slide_parts.extend(iter_pptx_shape_text(shape))
        if slide_parts:
            slides.append(f"Slide {index}\n" + "\n".join(slide_parts))

    text = "\n\n".join(slides).strip()
    return {"ok": bool(text), "text": text, "slide_count": len(presentation.slides)}


def extract_pptx_from_zip(file_bytes):
    slides = []
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as archive:
            slide_names = sorted(
                [name for name in archive.namelist() if re.match(r"^ppt/slides/slide\d+\.xml$", name)],
                key=lambda value: int(re.search(r"slide(\d+)\.xml$", value).group(1)),
            )
            for index, name in enumerate(slide_names, start=1):
                blocks = xml_text_blocks(archive.read(name))
                if blocks:
                    slides.append(f"Slide {index}\n" + "\n".join(blocks))
    except (zipfile.BadZipFile, ET.ParseError):
        return {"ok": False, "reason": "pptx_zip_parse_failed"}

    text = "\n\n".join(slides).strip()
    return {"ok": bool(text), "text": text, "slide_count": len(slides)}


def extract_pptx_text(file_bytes):
    extraction = extract_with_python_pptx(file_bytes)
    if extraction.get("ok"):
        extraction["mode"] = "python-pptx"
        return extraction

    extraction = extract_pptx_from_zip(file_bytes)
    if extraction.get("ok"):
        extraction["mode"] = "pptx-xml"
    return extraction


def extract_pdf_text(file_bytes, bucket, key):
    extraction_mode = "pdfplumber"
    extraction = extract_with_pdfplumber(file_bytes)

    if not extraction.get("ok"):
        extraction_mode = "pypdf"
        extraction = extract_with_pypdf(file_bytes)

    if (not extraction.get("ok") or extraction.get("scanned_like") or extraction.get("table_poor")) and USE_TEXTRACT_FALLBACK:
        extraction_mode = "textract"
        textract_result = extract_with_textract(bucket, key)
        extraction = {
            "ok": bool(textract_result.get("text")),
            "text": textract_result.get("text", ""),
            "textract_job_id": textract_result.get("job_id"),
        }

    extraction["mode"] = extraction_mode
    return extraction


def extract_document_text(file_bytes, filename, bucket, key):
    extension = PurePosixPath(filename).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        return {"ok": False, "reason": "unsupported_file_type", "file_type": extension or "unknown"}

    if extension == ".pdf":
        extraction = extract_pdf_text(file_bytes, bucket, key)
    elif extension in {".txt", ".md", ".markdown"}:
        extraction = extract_plain_text(file_bytes)
        extraction["mode"] = extension.lstrip(".")
    elif extension == ".vtt":
        extraction = extract_vtt_text(file_bytes)
    elif extension == ".docx":
        extraction = extract_docx_text(file_bytes)
    elif extension == ".pptx":
        extraction = extract_pptx_text(file_bytes)
    else:
        extraction = {"ok": False, "reason": "unsupported_file_type"}

    extraction["file_type"] = extension.lstrip(".")
    return extraction


def extract_with_textract(bucket, key):
    job = textract.start_document_text_detection(
        DocumentLocation={"S3Object": {"Bucket": bucket, "Name": key}}
    )
    job_id = job["JobId"]

    status = "IN_PROGRESS"
    attempts = 0
    while status == "IN_PROGRESS" and attempts < 90:
        attempts += 1
        time.sleep(2)
        poll = textract.get_document_text_detection(JobId=job_id, MaxResults=1000)
        status = poll.get("JobStatus", "IN_PROGRESS")
        if status in {"SUCCEEDED", "FAILED", "PARTIAL_SUCCESS"}:
            break

    if status not in {"SUCCEEDED", "PARTIAL_SUCCESS"}:
        raise RuntimeError(f"Textract job failed: {status}")

    lines = []
    token = None
    while True:
        page = textract.get_document_text_detection(JobId=job_id, MaxResults=1000, NextToken=token) if token else textract.get_document_text_detection(JobId=job_id, MaxResults=1000)
        for block in page.get("Blocks", []):
            if block.get("BlockType") == "LINE" and block.get("Text"):
                lines.append(block["Text"])
        token = page.get("NextToken")
        if not token:
            break
    return {"text": "\n".join(lines).strip(), "job_id": job_id}


def delete_processed_objects(bucket, user_id, session_id, doc_id):
    processed_root = KB_PROCESSED_PREFIX.rstrip("/")
    keys = [
        f"{processed_root}/{user_id}/{session_id}/{doc_id}.txt",
        f"documents/processed/{user_id}/{doc_id}.txt",
    ]
    chunk_prefixes = [
        f"{processed_root}/{user_id}/{session_id}/{doc_id}/",
        f"documents/processed/{user_id}/{doc_id}/",
    ]

    for chunk_prefix in chunk_prefixes:
        token = None
        while True:
            kwargs = {"Bucket": bucket, "Prefix": chunk_prefix}
            if token:
                kwargs["ContinuationToken"] = token
            response_data = s3.list_objects_v2(**kwargs)
            keys.extend([item["Key"] for item in response_data.get("Contents", [])])
            token = response_data.get("NextContinuationToken")
            if not token:
                break

    for index in range(0, len(keys), 1000):
        batch = list(dict.fromkeys(keys[index:index + 1000]))
        if batch:
            s3.delete_objects(Bucket=bucket, Delete={"Objects": [{"Key": key} for key in batch], "Quiet": True})


def upload_processed_text(bucket, user_id, session_id, doc_id, text):
    processed_root = KB_PROCESSED_PREFIX.rstrip("/")
    safe_session_id = session_id or "default"
    processed_prefix = f"{processed_root}/{user_id}/{safe_session_id}/"
    processed_key = f"{processed_prefix}{doc_id}.txt"
    delete_processed_objects(bucket, user_id, safe_session_id, doc_id)

    s3.put_object(
        Bucket=bucket,
        Key=processed_key,
        Body=text.encode("utf-8"),
        ContentType="text/plain; charset=utf-8",
    )

    return {
        "key": processed_key,
        "prefix": processed_prefix,
        "size_bytes": len(text.encode("utf-8")),
    }


def start_kb_ingestion():
    if not BEDROCK_KNOWLEDGE_BASE_ID or not BEDROCK_DATA_SOURCE_ID:
        raise RuntimeError("Missing BEDROCK_KNOWLEDGE_BASE_ID or BEDROCK_DATA_SOURCE_ID")
    response_data = bedrock_agent.start_ingestion_job(
        knowledgeBaseId=BEDROCK_KNOWLEDGE_BASE_ID,
        dataSourceId=BEDROCK_DATA_SOURCE_ID,
        clientToken=str(uuid.uuid4()),
        description="Document processing Lambda ingestion trigger",
    )
    return response_data.get("ingestionJob", {})


def latest_active_ingestion_job():
    response_data = bedrock_agent.list_ingestion_jobs(
        knowledgeBaseId=BEDROCK_KNOWLEDGE_BASE_ID,
        dataSourceId=BEDROCK_DATA_SOURCE_ID,
    )
    jobs = response_data.get("ingestionJobSummaries", [])
    def job_timestamp(job):
        value = job.get("updatedAt") or job.get("startedAt") or ""
        return value.isoformat() if hasattr(value, "isoformat") else str(value)

    active_jobs = [
        job for job in jobs
        if str(job.get("status", "")).upper() in {"STARTING", "IN_PROGRESS"}
    ]
    if active_jobs:
        return sorted(active_jobs, key=job_timestamp, reverse=True)[0]
    if jobs:
        return sorted(jobs, key=job_timestamp, reverse=True)[0]
    return {}


def is_ingestion_already_running_error(exc):
    if not isinstance(exc, ClientError):
        return False
    error = exc.response.get("Error", {})
    code = str(error.get("Code", ""))
    message = str(error.get("Message", "")).lower()
    if code in {"ConflictException", "Conflict"}:
        return True
    return "ingestion" in message and any(term in message for term in ["in progress", "running", "concurrent", "active"])


def process_record(record):
    detail = record.get("detail") or {}
    bucket = (detail.get("bucket") or {}).get("name")
    key = (detail.get("object") or {}).get("key")
    if not bucket or not key:
        # Support direct S3 event shape as fallback
        s3_record = (record.get("s3") or {})
        bucket = ((s3_record.get("bucket") or {}).get("name")) or bucket
        key = ((s3_record.get("object") or {}).get("key")) or key
    if not bucket or not key:
        return {"skipped": True, "reason": "missing_bucket_or_key"}

    key = unquote_plus(key)
    parsed = parse_object_key(key)
    if not parsed:
        return {"skipped": True, "reason": "key_not_in_documents_prefix", "key": key}

    user_id = parsed["user_id"]
    doc_id = parsed["doc_id"]
    session_id = parsed.get("session_id") or "default"
    filename = parsed["filename"]

    update_doc(
        user_id,
        doc_id,
        kb_status="PROCESSING",
        processing_started_at=now_iso(),
        title=filename,
        raw_s3_key=key,
        session_id=session_id,
    )

    extension = PurePosixPath(filename).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        update_doc(
            user_id,
            doc_id,
            kb_status="FAILED",
            failure_reason="unsupported_file_type",
            processed_at=now_iso(),
        )
        return {"doc_id": doc_id, "status": "FAILED", "reason": "unsupported_file_type", "file_type": extension or "unknown"}

    obj = s3.get_object(Bucket=bucket, Key=key)
    file_bytes = obj["Body"].read()
    extraction = extract_document_text(file_bytes, filename, bucket, key)
    extraction_mode = extraction.get("mode", "unknown")

    text = (extraction.get("text") or "").strip()
    if not text:
        reason = extraction.get("reason", "text_extraction_failed")
        update_doc(user_id, doc_id, kb_status="FAILED", failure_reason=reason, processed_at=now_iso())
        return {"doc_id": doc_id, "status": "FAILED", "reason": reason, "file_type": extraction.get("file_type", extension or "unknown")}

    processed_upload = upload_processed_text(bucket, user_id, session_id, doc_id, text)
    processed_key = processed_upload["key"]
    processed_attrs = {
        "extraction_mode": extraction_mode,
        "file_type": extraction.get("file_type", extension.lstrip(".")),
        "page_count": extraction.get("page_count", 0),
        "slide_count": extraction.get("slide_count", 0),
        "cue_count": extraction.get("cue_count", 0),
        "processed_text_size_bytes": processed_upload["size_bytes"],
        "processed_s3_prefix": processed_upload["prefix"],
        "raw_s3_key": key,
        "processed_s3_key": processed_key,
        "processed_text_s3_key": processed_key,
        "processed_at": now_iso(),
    }

    update_doc(user_id, doc_id, kb_status="PROCESSING", ingestion_status="STARTING", **processed_attrs)

    try:
        ingestion_job = start_kb_ingestion()
    except Exception as exc:
        if is_ingestion_already_running_error(exc):
            try:
                ingestion_job = latest_active_ingestion_job()
            except Exception:
                ingestion_job = {}
            ingestion_job_id = ingestion_job.get("ingestionJobId", "")
            ingestion_status = ingestion_job.get("status", "IN_PROGRESS")
            update_doc(
                user_id,
                doc_id,
                kb_status="PROCESSING",
                ingestion_job_id=ingestion_job_id,
                ingestion_status=ingestion_status,
                ingestion_note="Attached to active KB ingestion job after concurrent start was rejected.",
                **processed_attrs,
            )
            return {
                "doc_id": doc_id,
                "status": "PROCESSING",
                "reason": "ingestion_already_running",
                "processed_text_s3_key": processed_key,
                "ingestion_job_id": ingestion_job_id,
                "ingestion_status": ingestion_status,
            }
        update_doc(
            user_id,
            doc_id,
            kb_status="FAILED",
            ingestion_status="START_FAILED",
            failure_reason=f"ingestion_start_failed: {exc}",
            **processed_attrs,
        )
        return {
            "doc_id": doc_id,
            "status": "FAILED",
            "reason": "ingestion_start_failed",
            "processed_text_s3_key": processed_key,
            "error": str(exc),
        }

    update_doc(
        user_id,
        doc_id,
        kb_status="PROCESSING",
        ingestion_job_id=ingestion_job.get("ingestionJobId", ""),
        ingestion_status=ingestion_job.get("status", "STARTING"),
        **processed_attrs,
    )

    return {
        "doc_id": doc_id,
        "status": "PROCESSING",
        "extraction_mode": extraction_mode,
        "file_type": extraction.get("file_type", extension.lstrip(".")),
        "processed_text_s3_key": processed_key,
        "ingestion_job_id": ingestion_job.get("ingestionJobId", ""),
        "ingestion_status": ingestion_job.get("status", "STARTING"),
    }


def lambda_handler(event, _context):
    records = event.get("Records") or [event]
    results = []
    for record in records:
        try:
            results.append(process_record(record))
        except Exception as exc:  # pragma: no cover
            results.append({"status": "ERROR", "error": str(exc), "record": record})
    return {"results": results}
